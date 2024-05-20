import fitz
import textract
import pandas as pd
from pptx import Presentation
import os
import glob
from langchain.chains import RetrievalQA
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from langchain.callbacks.manager import CallbackManager
from langchain_community.llms import Ollama
from langchain_community.embeddings.ollama import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
import streamlit as st
import time
import base64

page_bg_img = """
<style>
[data-testid = "stAppViewContainer"] {
background-image: url("https://lifeboat.com/blog.images/ai-revolution-how-to-profit-from-the-next-big-technology.jpg");
backgroud-repeat:no-repeat;
background-size:cover;
background-attachment:local;

}

[data-testid= "stVerticalBlock"]{
color: gray;
}

[data-testid= "stHeader"]{
backgraund-color:rgba(0,0,0,0);
}

</style>
"""



def extract_text_from_pdf(file_path):
    loader = PyPDFLoader(file_path)
    text=loader.load()
    return text

def extract_text_from_powerpoint(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text
    return text

def extract_text_from_other_file_types(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return text
    
def process_document(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.pdf':
        text = extract_text_from_pdf(file_path)
    elif file_extension == '.pptx':
        text = extract_text_from_powerpoint(file_path)
    elif file_extension == '.txt':
        text = extract_text_from_other_file_types(file_path)

    if text:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,
            chunk_overlap=200,
            length_function=len
        )
        all_chunks = text_splitter.split_documents(text)

        st.session_state.vectorstore = Chroma.from_documents(
            documents=all_chunks,
            embedding=OllamaEmbeddings(model="mistral")
        )
        st.session_state.vectorstore.persist()
        return text
    else:
        st.error("Failed to process the document.")
        return "" 
    
def process_directory(directory):
    documents = {}
    for filename in glob.glob(f"{directory}/**/*", recursive=True):
        try:
            text = process_document(filename)
            documents[filename] = text
        except Exception as e:
            print(f"Error processing document {filename}: {e}")
    return documents


if not os.path.exists('files'):
    os.mkdir('files')

if not os.path.exists('chroma'):
    os.mkdir('chroma')

if 'template' not in st.session_state:
    st.session_state.template = """You are a knowledgeable chatbot, here to help with questions of the user. Your tone should be professional and informative.

    Context: {context}
    History: {history}

    User: {question}
    Chatbot:"""
if 'prompt' not in st.session_state:
    st.session_state.prompt = PromptTemplate(
        input_variables=["history", "context", "question"],
        template=st.session_state.template,
    )
if 'memory' not in st.session_state:
    st.session_state.memory = ConversationBufferMemory(
        memory_key="history",
        return_messages=True,
        input_key="question"
    )
if 'vectorstore' not in st.session_state:
    st.session_state.vectorstore = Chroma(persist_directory='local-rag',
                                          embedding_function=OllamaEmbeddings(model="nomic-embed-text")
                                          )
if 'llm' not in st.session_state:
    st.session_state.llm = Ollama(model="mistral",
                                  verbose=True,
                                  callback_manager=CallbackManager([StreamingStdOutCallbackHandler()])
                                  )



def load_image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode()
    
user_avatar_path = 'user.png'
assistant_avatar_path = 'assistente.png'

# Carregar os avatares em base64
user_avatar_base64 = load_image_to_base64(user_avatar_path)
assistant_avatar_base64 = load_image_to_base64(assistant_avatar_path)

# Initialize session state
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

#st.markdown(page_bg_img, unsafe_allow_html=True)    

st.title("PDF Assistant")

st.sidebar.image('logo.png', use_column_width=True)

st.markdown("""
    ## How to Use This Application
    - **Drag and drop** your PDF file into the designated area or use the upload button below.
    - Once you see a message stating your document has been processed, you can start asking questions in the chat input to interact with the PDF content.
    ---

""", unsafe_allow_html=True)

uploaded_file = st.sidebar.file_uploader("Upload your file", type=['pdf', 'txt', 'doc', 'docx', 'pptx'])
processar_doc = st.sidebar.button("Processar Documentos")

for message in st.session_state.chat_history:
    with st.chat_message(message["role"]):
        if message["role"] == "user":
            avatar=f"data:image/png;base64,{user_avatar_base64 }"
            st.write("User’s question 🗣️")
        elif message["role"] == "assistant":
            st.write("Assistant's response 🤖")
            avatar=f"data:image/png;base64,{assistant_avatar_base64}"
        st.markdown(message["message"])

if uploaded_file is not None:
    if not os.path.isfile("files/"+uploaded_file.name):
        with st.status("Processing your document..."):
            bytes_data = uploaded_file.read()
            f = open("files/"+uploaded_file.name, "wb")
            f.write(bytes_data)

            processed_text = process_document("files/"+uploaded_file.name)
            if processed_text:
                st.success("Your document has been processed.")
                st.balloons()
                
    st.session_state.retriever = st.session_state.vectorstore.as_retriever()
    if 'qa_chain' not in st.session_state:
        st.session_state.qa_chain = RetrievalQA.from_chain_type(
            llm=st.session_state.llm,
            chain_type='stuff',
            retriever=st.session_state.retriever,
            verbose=True,
            chain_type_kwargs={
                "verbose": True,
                "prompt": st.session_state.prompt,
                "memory": st.session_state.memory,
            }
        )
    
    if user_input := st.chat_input("You:", key="user_input"):
        user_message = {"role": "user", "message": user_input}
        st.session_state.chat_history.append(user_message)
        with st.chat_message("user", avatar=f"data:image/png;base64,{user_avatar_base64 }"):
            st.write("User’s question 🗣️")
            st.markdown(user_input)
        
        with st.chat_message("assistant", avatar=f"data:image/png;base64,{assistant_avatar_base64}"):
            st.write("Assistant's response 🤖")
            with st.spinner("Assistant is typing..."):
                response = st.session_state.qa_chain(user_input)
            message_placeholder = st.empty()
            full_response = ""
            for chunk in response['result'].split():
                full_response += chunk + " "
                time.sleep(0.05)
                message_placeholder.markdown(full_response + "▌")
            message_placeholder.markdown(full_response)

        chatbot_message = {"role": "assistant", "message": response['result']}
        st.session_state.chat_history.append(chatbot_message)
        
else:
    st.warning("Please upload a file and click the process button.", icon="⚠️")
